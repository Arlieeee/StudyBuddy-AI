"""Document parser for PDF, PPTX, DOCX, and TXT files."""
import os
from pathlib import Path
from typing import List, Tuple, Optional
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document

from ..config import get_settings


class DocumentParser:
    """Parse various document formats and extract text."""
    
    def __init__(self):
        """Initialize document parser."""
        self.settings = get_settings()
    
    def parse(self, file_path: str) -> Tuple[str, dict]:
        """
        Parse document and extract text.
        
        Args:
            file_path: Path to the document file
        
        Returns:
            Tuple of (extracted_text, metadata)
        """
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        parsers = {
            ".pdf": self._parse_pdf,
            ".pptx": self._parse_pptx,
            ".docx": self._parse_docx,
            ".txt": self._parse_txt,
        }
        
        parser = parsers.get(extension)
        if not parser:
            raise ValueError(f"Unsupported file format: {extension}")
        
        text, metadata = parser(str(file_path))
        metadata["filename"] = file_path.name
        metadata["file_type"] = extension[1:]  # Remove the dot
        
        return text, metadata
    
    def _parse_pdf(self, file_path: str) -> Tuple[str, dict]:
        """Parse PDF file."""
        doc = fitz.open(file_path)
        text_parts = []
        page_count = len(doc)  # Save page count before closing
        
        for page_num, page in enumerate(doc, 1):
            page_text = page.get_text()
            if page_text.strip():
                text_parts.append(f"[Page {page_num}]\n{page_text}")
        
        doc.close()
        
        metadata = {
            "page_count": page_count,
            "format": "pdf",
        }
        
        return "\n\n".join(text_parts), metadata
    
    def _parse_pptx(self, file_path: str) -> Tuple[str, dict]:
        """Parse PowerPoint file."""
        prs = Presentation(file_path)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                text_parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))
        
        metadata = {
            "slide_count": len(prs.slides),
            "format": "pptx",
        }
        
        return "\n\n".join(text_parts), metadata
    
    def _parse_docx(self, file_path: str) -> Tuple[str, dict]:
        """Parse Word document."""
        doc = Document(file_path)
        text_parts = []
        
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text for cell in row.cells if cell.text.strip()]
                if row_text:
                    text_parts.append(" | ".join(row_text))
        
        metadata = {
            "paragraph_count": len(doc.paragraphs),
            "format": "docx",
        }
        
        return "\n\n".join(text_parts), metadata
    
    def _parse_txt(self, file_path: str) -> Tuple[str, dict]:
        """Parse plain text file."""
        encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
        
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    text = f.read()
                break
            except UnicodeDecodeError:
                continue
        else:
            raise ValueError("Unable to decode text file")
        
        metadata = {
            "char_count": len(text),
            "format": "txt",
        }
        
        return text, metadata
    
    def chunk_text(
        self,
        text: str,
        chunk_size: Optional[int] = None,
        chunk_overlap: Optional[int] = None,
    ) -> List[str]:
        """
        Split text into overlapping chunks.
        
        Args:
            text: Text to chunk
            chunk_size: Size of each chunk
            chunk_overlap: Overlap between chunks
        
        Returns:
            List of text chunks
        """
        chunk_size = chunk_size or self.settings.CHUNK_SIZE
        chunk_overlap = chunk_overlap or self.settings.CHUNK_OVERLAP
        
        if len(text) <= chunk_size:
            return [text]
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            
            # Try to break at sentence boundaries
            if end < len(text):
                # Look for sentence endings
                for sep in ["。", ".", "!\n", "?\n", "\n\n"]:
                    last_sep = text.rfind(sep, start, end)
                    if last_sep > start + chunk_size // 2:
                        end = last_sep + len(sep)
                        break
            
            chunk = text[start:end].strip()
            if chunk:
                chunks.append(chunk)
            
            start = end - chunk_overlap
            if start >= len(text):
                break
        
        return chunks


# Global parser instance
_parser: Optional[DocumentParser] = None


def get_document_parser() -> DocumentParser:
    """Get or create document parser instance."""
    global _parser
    if _parser is None:
        _parser = DocumentParser()
    return _parser
